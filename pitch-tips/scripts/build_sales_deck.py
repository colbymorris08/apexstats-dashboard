#!/usr/bin/env python3
"""Rebuild Preflight sales deck (markdown, HTML, PPTX, PDF) for Downloads."""
from __future__ import annotations

import re
import shutil
import subprocess
import textwrap
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pypdf import PdfReader, PdfWriter

OUT = Path("/Users/colbymorris/Downloads/Preflight_Sales_Deck")
DELIVERABLES = Path("/Users/colbymorris/Downloads/Preflight_Deliverables/05_slide_deck")
PACKAGE_DATE = "September 2, 2026"

SLIDES = [
    {
        "title": "Preflight: Pre-Release Pitch Intelligence from Broadcast Video",
        "subhead": "Colby Morris · Computer Vision Scouting Platform",
        "bullets": [
            "Identify pitch type **1.4+ seconds before release** from standard broadcast footage",
            "No radar, no wearable sensors, no proprietary camera installs",
            "Validated across MLB, NPB, KBO, CPBL, NCAA, MiLB, and LMB",
            "Every arm profiled across **all pitch types × all situation filters**",
        ],
        "extra_html": """
<div class="brand-mark">PREFLIGHT</div>
<p class="tagline">Computer Vision Scouting Platform</p>""",
    },
    {
        "title": "Hitters Have ~400ms. Scouting Has Almost None.",
        "subhead": "Post-release identification is too late for swing decisions.",
        "bullets": [
            "Traditional pitch ID relies on post-release movement — too late for swing decisions",
            "Human scouts catch ~60–70% of mechanical tells under live game speed",
            "International and amateur arms have zero systematic pre-release scouting coverage",
            "Existing tools (TrackMan, Hawkeye) measure **outcome**, not **intent**",
        ],
        "extra_html": """
<div class="timeline" aria-label="Pitch delivery timeline">
  <div class="tl-step">Set Position</div><div class="tl-arrow">→</div>
  <div class="tl-step">Leg Lift</div><div class="tl-arrow">→</div>
  <div class="tl-step">Hand Break</div><div class="tl-arrow">→</div>
  <div class="tl-step highlight">Release</div><div class="tl-arrow">→</div>
  <div class="tl-step highlight-warn">Plate <span>~400ms window</span></div>
</div>""",
    },
    {
        "title": "Physical Mechanical Separation, Measured Before Hand Break",
        "subhead": "Computer vision isolates actionable pre-release mechanical tells from any broadcast CF angle.",
        "bullets": [
            "Computer vision isolates **glove set height, hand depth, tempo cadence, and torso spacing** in the pre-release window",
            "Each indicator carries an empirical discrimination score (Signal Floor ≥75%)",
            "Delivered as ranked scouting dossiers with **independent delivery compare scrubbers**",
            "Works on any standard broadcast center-field camera angle",
        ],
        "images": [("roupp_sinker_vs_curve_deck_comparison.png", "Roupp Sinker vs Curve comparison", "hero")],
    },
    {
        "title": "Case Study: 88.4% Signal on Glove Set Height Discrimination",
        "subhead": "Glove set height SI vs CU — +4.8 in elevation delta",
        "bullets": [
            "**Primary tell:** Glove rests at belt buckle (SI) vs chest lettering (CU) — +4.8 in elevation delta",
            "**Secondary tell:** Wrist sits 1.6 in deeper in glove web pocket on changeup grip",
            "**Tempo tell:** +140ms longer glove pause before leg lift on breaking pitches",
            "**Actionability:** Opposition hitters can identify pitch group at first baseline touch, **1.42s before release**",
            "**Validation:** Out-of-sample holdout accuracy 89%, effect size d = 1.42, 5.9× Signal Floor",
        ],
        "images": [
            ("roupp_sinker_vs_curve_deck_comparison.png", "Roupp comparison card", "hero"),
            ("roupp_si_sinker_still.png", "SI still", "thumb"),
            ("roupp_cu_curve_still.png", "CU still", "thumb"),
            ("roupp_si_sinker_set.png", "SI set", "thumb"),
            ("roupp_cu_curve_set.png", "CU set", "thumb"),
        ],
    },
    {
        "title": "Catcher Pre-Pitch Target Discrimination: 91.2%",
        "subhead": "FF vs CH target — +6.2 in elevation delta",
        "bullets": [
            "Fastball target held **6.2 inches higher** than offspeed/breaker presentation",
            "One-knee-down block posture committed on offspeed; active crouch on fastballs",
            "Glove face angled upward on heaters, turned inward on changeup receiving",
            "Signal available **2.1 seconds** prior to pitch release during sign-calling sequence",
            "Dugout & base runners can relay pitch group confirmation before pitcher motion onset",
        ],
        "images": [
            ("moreno_catcher_setup_deck_comparison.png", "Moreno FF vs CH comparison", "hero"),
            ("moreno_ff_fastball_setup.png", "FF setup", "thumb"),
            ("moreno_ch_changeup_setup.png", "CH setup", "thumb"),
        ],
    },
    {
        "title": "The Preflight Publishing Workflow",
        "subhead": "Per player: all pitches × all situation filters — quota-sampled, verified, and hitter-ready.",
        "bullets": [
            "**Coverage grid:** Every arm is profiled across **all pitch types × all situation filters** (bases empty, runners on, runner on 2nd, two strikes, and more)",
            "**Quota sampling:** **n≈10 pitches per pitch type per situation bucket** — newest games first, capped per game so no single outing dominates",
            "**Exemplar clips:** **One published MP4 per situation** with full tip metadata — discrimination score, cue window, and apex timestamp",
            "**Hitter prep surface:** Player pages show the **most recent verified example + game date** so advance staff know the film is current",
            "**Trust layer:** **Identity verification** on every clip, **independent per-pane scrubbers** with **apex labels**, MD5 integrity checks — **no imposters**",
        ],
        "extra_html": """
<div class="workflow-flow">
  <span class="pipe-node">All Pitches</span>
  <span class="pipe-node">×</span>
  <span class="pipe-node">All Situations</span>
  <span class="pipe-node">→</span>
  <span class="pipe-node">n≈10 / cell</span>
  <span class="pipe-node">→</span>
  <span class="pipe-node">Exemplar MP4</span>
  <span class="pipe-node">→</span>
  <span class="pipe-node">Hitter Prep</span>
</div>
<div class="img-grid small">
<figure><img src="sample_roupp_si_f132.jpg" alt="Roupp SI frame"><figcaption>Roupp SI apex frame</figcaption></figure>
<figure><img src="sample_moreno_ff_f096.jpg" alt="Moreno FF frame"><figcaption>Moreno FF setup frame</figcaption></figure>
</div>""",
    },
    {
        "title": "The Preflight CV Pipeline",
        "subhead": "Six validated layers — YOLOv8 fine-tuned across six international league domains.",
        "bullets": [
            "**Pose + detection:** MediaPipe 33-point tracking + YOLOv8m fine-tuned on 12,000+ broadcast frames (91.2% mAP@0.5 glove detection)",
            "**Normalization:** Torso-relative units eliminate camera distance, zoom, and player height confounders",
            "**Zero-leakage window:** Features anchored set position through peak leg lift, ending before hand break (1.2–2.1s actionability)",
            "**Holdout validation:** 80/20 temporal split — reported accuracy is holdout-only, not in-sample inflated",
            "**Signal Floor ≥75%:** Indicators below threshold are rejected — no noise enters the scouting dossier",
            "**Transfer learning:** MLB-trained weights transfer to international leagues with <500 additional annotated frames per league",
        ],
        "extra_html": """
<div class="pipeline-flow">
<span class="pipe-node">Video Input</span>
<span class="pipe-node">Pose Tracking</span>
<span class="pipe-node">Object Detection</span>
<span class="pipe-node">Feature Norm</span>
<span class="pipe-node">Discrimination</span>
<span class="pipe-node">Scouting Output</span>
</div>
<div class="img-grid small svg-grid">
<figure><img src="detection/sf_landen_roupp_si_f132.svg" alt="Roupp SI detection"><figcaption>Roupp SI detection</figcaption></figure>
<figure><img src="detection/ari_gabriel_moreno_ff_f096.svg" alt="Moreno FF detection"><figcaption>Moreno FF detection</figcaption></figure>
</div>""",
    },
    {
        "title": "91 Arms Modeled · Full Situational Grid · 227 Clips Verified",
        "subhead": "All pitch types × all situation filters — quota-sampled and integrity-checked.",
        "bullets": [
            "**MLB:** Full NL West coverage (ARI, COL, LAD, SD, SF) — 28 pitchers, 5 catchers",
            "**International:** NPB, KBO, CPBL, LMB, NCAA, MiLB pipelines active",
            "**Per-arm grid:** Every pitch type crossed with every situation filter — **n≈10 samples per cell**",
            "**227 verified non-MLB video clips** with MD5 integrity checks (see non_mlb_video_verification.json)",
            "**One exemplar MP4 per situation** published with tip metadata; **most recent example + date** on every player page",
            "Enterprise pilot unlocks full dossier access for any arm in the database",
        ],
        "extra_html": """
<div class="coverage-grid">
  <div class="stat-card"><div class="stat-num">91</div><div class="stat-label">Arms Modeled</div></div>
  <div class="stat-card"><div class="stat-num">4+</div><div class="stat-label">Situation Filters</div></div>
  <div class="stat-card"><div class="stat-num">~10</div><div class="stat-label">Pitches / Cell</div></div>
  <div class="stat-card"><div class="stat-num">227</div><div class="stat-label">Clips Verified</div></div>
</div>
<div class="league-pins">MLB · NPB · KBO · CPBL · NCAA · MiLB · LMB</div>""",
    },
    {
        "title": 'What "2 of 5" Means — and Why It Matters',
        "subhead": "Only statistically distinct pre-release mechanics enter the scouting dossier.",
        "bullets": [
            'A 5-pitch arsenal (FF, SL, CH, SI, CU) is tested for mechanical separation',
            '"2 of 5" = 2 pitch types exhibit statistically distinct pre-release mechanics (≥75% discrimination)',
            "The other 3 pitches share delivery signatures — no actionable tell exists for those pairs",
            "Preflight only surfaces indicators that clear the Signal Floor — zero false-positive scouting noise",
            "Example: Roupp's SI and CU are discriminable (2 of 5); his CH shares SI's glove set (not separately identifiable)",
        ],
        "extra_html": """
<div class="signal-card">
  <div class="signal-badge">Signal Floor ≥ 75%</div>
  <div class="signal-example"><span class="on">SI</span><span class="on">CU</span><span class="off">CH</span><span class="off">SL</span><span class="off">FF</span></div>
  <p class="signal-caption">2 of 5 pitch types show distinct pre-release mechanics</p>
</div>""",
    },
    {
        "title": "Request Enterprise Scouting Pilot Access",
        "subhead": "Verified clips, independent scrubbers, and full situational dossiers.",
        "bullets": [
            "**Lite (Public):** 3 fully unlocked showcase arms (Rodriguez, Webb, Roupp) + catcher demo (Moreno)",
            "**Enterprise Pilot:** Full dossier access for any arm in the database",
            "Independent delivery compare scrubbers with per-clip **apex labels**",
            "Situational breakdown tables — all pitch types × all situation filters",
            "**Identity verification** on every published clip — MD5-checked, **no imposters**",
            "Out-of-sample holdout validation reports",
            "**Contact:** @colbymorris08 on X · Enterprise pilot request form on platform",
        ],
        "extra_html": """
<div class="tier-grid">
  <div class="tier lite"><h3>Lite (Public)</h3><ul><li>3 showcase pitchers unlocked</li><li>Catcher demo (Moreno)</li><li>preflightpitchtips.com</li></ul></div>
  <div class="tier enterprise"><h3>Enterprise Pilot</h3><ul><li>Full database access</li><li>Independent scrubbers + apex labels</li><li>Holdout validation reports</li></ul></div>
</div>""",
    },
    {
        "title": "See the Signal Before the Release",
        "subhead": "Live demo available now — enterprise pilots open for full roster access.",
        "bullets": [
            "Live demo available at preflightpitchtips.com",
            "Showcase arms fully unlocked: Eduardo Rodriguez (ARI), Logan Webb (SF), Landen Roupp (SF)",
            "Catcher intelligence demo: Gabriel Moreno (AZ)",
            "Every player page: **most recent verified example + game date** for hitter prep",
            "Enterprise pilot inquiries: DM @colbymorris08",
        ],
        "images": [
            ("roupp_sinker_vs_curve_deck_comparison.png", "Roupp comparison", "duo"),
            ("moreno_catcher_setup_deck_comparison.png", "Moreno comparison", "duo"),
        ],
    },
]

assert len(SLIDES) == 11, f"Expected 11 slides, got {len(SLIDES)}"


def md_bold(text: str) -> str:
    return re.sub(r"\*\*(.+?)\*\*", r"**\1**", text)


def html_bold(text: str) -> str:
    return re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)


def render_images_html(images: list[tuple[str, str, str]]) -> str:
    if not images:
        return ""
    heroes = [i for i in images if i[2] == "hero"]
    thumbs = [i for i in images if i[2] == "thumb"]
    duos = [i for i in images if i[2] == "duo"]
    parts = []
    if heroes:
        parts.append('<div class="img-grid">')
        for src, alt, _ in heroes:
            parts.append(f'<figure class="hero-img"><img src="{src}" alt="{alt}"></figure>')
        parts.append("</div>")
    if thumbs:
        parts.append('<div class="img-grid thumbs">')
        for src, alt, _ in thumbs:
            parts.append(f"<figure><img src=\"{src}\" alt=\"{alt}\"><figcaption>{alt}</figcaption></figure>")
        parts.append("</div>")
    if duos:
        parts.append('<div class="img-grid duo">')
        for src, alt, _ in duos:
            parts.append(f'<figure><img src="{src}" alt="{alt}"></figure>')
        parts.append("</div>")
    return "\n".join(parts)


def build_markdown() -> str:
    lines = [
        "# Preflight Sales Pitch Deck — Full Slide Content",
        "",
        f"**Author:** Colby Morris · Preflight Computer Vision Scouting Platform  ",
        f"**Package Date:** {PACKAGE_DATE}  ",
        "**Assets folder:** Same directory as this document (`Preflight_Sales_Deck/`)",
        "",
        "---",
        "",
    ]
    for i, slide in enumerate(SLIDES, 1):
        lines.append(f"## SLIDE {i} — {slide['title'].split(':')[0] if ':' in slide['title'] else slide['title'][:40]}")
        lines.append("")
        lines.append(f"**Headline:** {slide['title']}")
        lines.append("")
        lines.append(f"**Subhead:** {slide['subhead']}")
        lines.append("")
        if slide.get("images"):
            lines.append("**Recommended image:** " + ", ".join(f"`{src}`" for src, _, _ in slide["images"]))
            lines.append("")
        lines.append("**Bullet points:**")
        for b in slide["bullets"]:
            lines.append(f"- {b}")
        lines.append("")
        lines.append("---")
        lines.append("")
    lines.extend([
        "## Image Quick Reference",
        "",
        "| Slide | Primary Image File |",
        "|---|---|",
        "| 1 | Logo / brand graphic (no file in package) |",
        "| 2 | Timeline graphic (built-in CSS) |",
        "| 3 | `roupp_sinker_vs_curve_deck_comparison.png` |",
        "| 4 | `roupp_sinker_vs_curve_deck_comparison.png` |",
        "| 5 | `moreno_catcher_setup_deck_comparison.png` |",
        "| 6 | Workflow diagram + `sample_roupp_si_f132.jpg` |",
        "| 7 | CV pipeline + `detection/*.svg` overlays |",
        "| 8 | Platform stats + `non_mlb_video_verification.json` |",
        "| 9 | Signal Floor explainer card |",
        "| 10 | Lite vs Enterprise tier cards |",
        "| 11 | Both comparison PNGs side-by-side |",
    ])
    return "\n".join(lines) + "\n"


HTML_HEAD = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Preflight Sales Deck — 11 Slides</title>
<style>
:root {
  --bg: #070b14; --bg2: #0f1729; --text: #e8edf5; --muted: #94a3b8;
  --accent: #22d3a8; --accent2: #38bdf8; --warn: #f59e0b;
  --card: #111827; --border: #1e293b;
}
* { box-sizing: border-box; margin: 0; padding: 0; }
html, body { height: 100%; background: var(--bg); color: var(--text);
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif; overflow: hidden; }
.deck { height: 100vh; position: relative; }
.slide {
  display: none; height: 100vh; padding: 2.5rem 3.5rem 4.5rem; overflow-y: auto;
  background: radial-gradient(ellipse at 20% 0%, #132035 0%, var(--bg) 55%);
}
.slide.active { display: block; }
.slide-num { position: fixed; bottom: 1.2rem; right: 2rem; color: var(--muted); font-size: 0.85rem; z-index: 5; }
h1 { font-size: clamp(1.6rem, 3.2vw, 2.4rem); line-height: 1.15; max-width: 52rem; margin-bottom: 0.6rem; }
.subhead { color: var(--muted); font-size: clamp(1rem, 1.6vw, 1.15rem); margin-bottom: 1.2rem; max-width: 48rem; }
.bullets { margin-top: 1rem; padding-left: 1.4rem; max-width: 52rem; }
.bullets li { margin: 0.45rem 0; line-height: 1.45; font-size: clamp(0.92rem, 1.3vw, 1.05rem); }
.bullets strong { color: var(--accent); }
.brand-mark { font-size: 3rem; font-weight: 800; letter-spacing: 0.25em; color: var(--accent); margin: 2rem 0 0.5rem; }
.tagline { color: var(--muted); font-size: 1.1rem; margin-bottom: 2rem; }
.timeline { display: flex; flex-wrap: wrap; align-items: center; gap: 0.4rem; margin: 1.5rem 0; padding: 1rem; background: var(--card); border: 1px solid var(--border); border-radius: 12px; max-width: 52rem; }
.tl-step { padding: 0.5rem 0.8rem; background: var(--bg2); border-radius: 8px; font-size: 0.9rem; }
.tl-step.highlight { background: #14532d; color: var(--accent); font-weight: 600; }
.tl-step.highlight-warn { background: #78350f; color: var(--warn); font-weight: 600; }
.tl-step span { display: block; font-size: 0.75rem; opacity: 0.9; }
.tl-arrow { color: var(--muted); }
.img-grid { display: grid; grid-template-columns: 1fr; gap: 0.8rem; margin: 1rem 0; max-width: 56rem; }
.img-grid.duo { grid-template-columns: 1fr 1fr; }
.img-grid.thumbs { grid-template-columns: repeat(auto-fit, minmax(140px, 1fr)); }
.img-grid.small { grid-template-columns: 1fr 1fr; max-width: 36rem; }
.img-grid img { width: 100%; height: auto; border-radius: 10px; border: 1px solid var(--border); }
.img-grid figcaption { font-size: 0.75rem; color: var(--muted); margin-top: 0.25rem; text-align: center; }
.hero-img img { max-height: 42vh; object-fit: contain; }
.pipeline-flow, .workflow-flow { display: flex; flex-wrap: wrap; gap: 0.35rem; margin: 0.8rem 0 1rem; max-width: 56rem; }
.pipe-node { background: var(--card); border: 1px solid var(--accent2); color: var(--accent2); padding: 0.35rem 0.6rem; border-radius: 999px; font-size: 0.78rem; }
.layers { display: grid; grid-template-columns: 1fr 1fr; gap: 0.6rem; max-width: 56rem; margin-top: 0.5rem; }
.layer { background: var(--card); border: 1px solid var(--border); border-radius: 8px; padding: 0.55rem 0.7rem; }
.layer h3 { font-size: 0.78rem; color: var(--accent2); margin-bottom: 0.3rem; }
.layer ul { padding-left: 1rem; }
.layer li { font-size: 0.72rem; line-height: 1.35; margin: 0.15rem 0; color: var(--muted); }
.data-table { width: 100%; max-width: 56rem; border-collapse: collapse; margin: 0.8rem 0; font-size: 0.82rem; }
.data-table th, .data-table td { border: 1px solid var(--border); padding: 0.4rem 0.55rem; text-align: left; }
.data-table th { background: var(--card); color: var(--accent2); }
.coverage-grid { display: grid; grid-template-columns: repeat(4, 1fr); gap: 0.8rem; max-width: 52rem; margin: 1.2rem 0; }
.stat-card { background: var(--card); border: 1px solid var(--border); border-radius: 12px; padding: 1rem; text-align: center; }
.stat-num { font-size: 2rem; font-weight: 700; color: var(--accent); }
.stat-label { font-size: 0.8rem; color: var(--muted); margin-top: 0.25rem; }
.league-pins { font-size: 0.95rem; color: var(--accent2); letter-spacing: 0.05em; margin-bottom: 0.5rem; }
.signal-card { background: var(--card); border: 1px solid var(--border); border-radius: 12px; padding: 1.2rem; max-width: 28rem; margin: 1rem 0; }
.signal-badge { display: inline-block; background: #14532d; color: var(--accent); padding: 0.3rem 0.7rem; border-radius: 999px; font-size: 0.8rem; font-weight: 600; margin-bottom: 0.8rem; }
.signal-example { display: flex; gap: 0.5rem; margin-bottom: 0.5rem; }
.signal-example span { width: 2.5rem; height: 2.5rem; border-radius: 8px; display: flex; align-items: center; justify-content: center; font-weight: 700; font-size: 0.85rem; }
.signal-example .on { background: #14532d; color: var(--accent); border: 2px solid var(--accent); }
.signal-example .off { background: var(--bg2); color: var(--muted); border: 2px solid var(--border); }
.signal-caption { color: var(--muted); font-size: 0.9rem; }
.tier-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 1rem; max-width: 44rem; margin: 1rem 0; }
.tier { background: var(--card); border: 1px solid var(--border); border-radius: 12px; padding: 1rem; }
.tier.enterprise { border-color: var(--accent); }
.tier h3 { font-size: 1rem; margin-bottom: 0.5rem; color: var(--accent2); }
.tier ul { padding-left: 1.1rem; }
.tier li { font-size: 0.88rem; margin: 0.3rem 0; color: var(--muted); }
.nav { position: fixed; bottom: 1rem; left: 50%; transform: translateX(-50%); display: flex; gap: 0.5rem; z-index: 10; }
.nav button { background: var(--card); color: var(--text); border: 1px solid var(--border); padding: 0.55rem 1.1rem; border-radius: 999px; cursor: pointer; font-size: 0.9rem; }
.nav button:hover { border-color: var(--accent); color: var(--accent); }
.help { position: fixed; bottom: 1.2rem; left: 2rem; color: var(--muted); font-size: 0.78rem; z-index: 5; }
.svg-grid img { background: #000; max-height: 18vh; object-fit: contain; }
@media (max-width: 900px) {
  .slide { padding: 1.5rem 1.2rem 4rem; }
  .layers, .tier-grid, .coverage-grid, .img-grid.duo, .img-grid.small { grid-template-columns: 1fr; }
}
@media print {
  html, body { overflow: visible; height: auto; }
  .slide { display: block !important; page-break-after: always; height: auto; min-height: 100vh; }
  .nav, .help { display: none; }
}
</style>
</head>
<body>
<div class="deck" id="deck">
"""


HTML_TAIL = """
</div>
<div class="help">← → arrow keys · click buttons to navigate</div>
<div class="nav">
  <button type="button" id="prev" aria-label="Previous slide">← Prev</button>
  <button type="button" id="next" aria-label="Next slide">Next →</button>
</div>
<script>
(function() {
  const slides = Array.from(document.querySelectorAll('.slide'));
  let idx = 0;
  function show(i) {
    idx = Math.max(0, Math.min(slides.length - 1, i));
    slides.forEach((s, n) => s.classList.toggle('active', n === idx));
    history.replaceState(null, '', '#slide-' + (idx + 1));
  }
  document.getElementById('prev').onclick = () => show(idx - 1);
  document.getElementById('next').onclick = () => show(idx + 1);
  document.addEventListener('keydown', (e) => {
    if (e.key === 'ArrowRight' || e.key === 'PageDown' || e.key === ' ') { e.preventDefault(); show(idx + 1); }
    if (e.key === 'ArrowLeft' || e.key === 'PageUp') { e.preventDefault(); show(idx - 1); }
    if (e.key === 'Home') show(0);
    if (e.key === 'End') show(slides.length - 1);
  });
  const hash = location.hash.match(/slide-(\\d+)/);
  if (hash) show(parseInt(hash[1], 10) - 1);
  else show(0);
})();
</script>
</body>
</html>
"""


def build_html() -> str:
    parts = [HTML_HEAD]
    for i, slide in enumerate(SLIDES, 1):
        parts.append(f'<section class="slide" data-slide="{i}" id="slide-{i}">')
        parts.append(f'<div class="slide-num">{i} / 11</div>')
        parts.append(f"<h1>{html_bold(slide['title'])}</h1>")
        parts.append(f"<p class=\"subhead\">{html_bold(slide['subhead'])}</p>")
        if slide.get("extra_html"):
            parts.append(slide["extra_html"])
        parts.append(render_images_html(slide.get("images", [])))
        if slide["bullets"]:
            parts.append("<ul class=\"bullets\">")
            for b in slide["bullets"]:
                parts.append(f"<li>{html_bold(b)}</li>")
            parts.append("</ul>")
        parts.append("</section>")
    parts.append(HTML_TAIL)
    return "\n".join(parts)


def build_readme() -> str:
    return textwrap.dedent(f"""\
        PREFLIGHT SALES DECK — OPEN-AND-USE DELIVERABLES
        =================================================
        Generated: {PACKAGE_DATE}
        Location: Preflight_Sales_Deck/

        These files work by double-click — no R, Node, or scripts required.


        PRIMARY: HTML SLIDE DECK (RECOMMENDED)
        --------------------------------------
        File: Preflight_Sales_Deck.html

        How to open:
          • Double-click Preflight_Sales_Deck.html (opens in Safari, Chrome, or Firefox)
          • Use ← → arrow keys, Space, or the Prev/Next buttons to navigate
          • 11 slides total — slide counter shown bottom-right

        How to present fullscreen:
          • Chrome/Safari: View → Enter Full Screen (or F11 in Chrome)
          • Press F5 or reload to restart at slide 1


        FULL PDF DECK
        ---------------
        File: Preflight_Sales_Deck.pdf

        How to open:
          • Double-click Preflight_Sales_Deck.pdf (Preview, Acrobat, or browser)
          • 11 pages — one slide per page, dark theme preserved

        Individual slide PDFs:
          • slides_pdf/slide_01.pdf through slide_11.pdf
          • Useful for email attachments or Google Slides import one-at-a-time


        OPTIONAL: POWERPOINT
        --------------------
        File: Preflight_Sales_Deck.pptx

        Auto-generated from slide content via python-pptx.
        Contains all slide copy. Best used as an editable starting point —
        import images from this folder manually, or use the HTML/PDF deck
        for a polished presentation.

        How to open:
          • Double-click in Keynote, PowerPoint, or Google Slides (File → Import)


        SOURCE COPY
        -----------
        Full markdown source with all bullet text:
          Preflight_Sales_Deck_Slides.md


        NEW WORKFLOW (Slide 6)
        ----------------------
        • Per player: all pitches × all situation filters
        • Sample n≈10 per pitch type per situation bucket
        • One exemplar MP4 per situation, published with tip metadata
        • Most recent example + date on site for hitter prep
        • Identity verification, independent scrub + apex labels, no imposters


        SLIDE-TO-IMAGE MAPPING
        ----------------------
        Slide  1 — Title (brand styling, no image)
        Slide  2 — Timeline graphic (built-in CSS)
        Slide  3 — roupp_sinker_vs_curve_deck_comparison.png
        Slide  4 — Roupp comparison + 4 supporting stills
        Slide  5 — moreno_catcher_setup_deck_comparison.png + setup stills
        Slide  6 — Publishing workflow diagram + sample JPG frames
        Slide  7 — CV pipeline diagram + detection SVGs
        Slide  8 — Coverage stats cards
        Slide  9 — Signal Floor explainer card
        Slide 10 — Lite vs Enterprise tier cards
        Slide 11 — Both comparison PNGs side-by-side


        GOOGLE SLIDES / KEYNOTE
        -----------------------
        Option A (fastest): Present directly from Preflight_Sales_Deck.html
        Option B: Open Preflight_Sales_Deck.pdf (full deck or slides_pdf/)
        Option C: Open Preflight_Sales_Deck.pptx, drag PNGs from this folder


        Contact: @colbymorris08 on X · preflightpitchtips.com
        """)


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, slide in enumerate(SLIDES, 1):
        s = prs.slides.add_slide(blank)
        # dark background
        bg = s.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(7, 11, 20)

        num_box = s.shapes.add_textbox(Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.4))
        num_tf = num_box.text_frame
        num_tf.text = f"{i} / 11"
        num_p = num_tf.paragraphs[0]
        num_p.font.size = Pt(11)
        num_p.font.color.rgb = RGBColor(148, 163, 184)

        title_box = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.text = slide["title"]
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(232, 237, 245)

        sub_box = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(0.6))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_tf.text = slide["subhead"]
        sub_p = sub_tf.paragraphs[0]
        sub_p.font.size = Pt(16)
        sub_p.font.color.rgb = RGBColor(148, 163, 184)

        body_box = s.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(12), Inches(4.5))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        for j, bullet in enumerate(slide["bullets"]):
            text = re.sub(r"\*\*(.+?)\*\*", r"\1", bullet)
            p = body_tf.paragraphs[0] if j == 0 else body_tf.add_paragraph()
            p.text = f"• {text}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(232, 237, 245)
            p.space_after = Pt(6)
            p.level = 0

    prs.save(path)


def build_pdf_from_html(html_path: Path, pdf_path: Path, slides_pdf_dir: Path) -> int:
    from io import BytesIO

    from PIL import Image
    from playwright.sync_api import sync_playwright
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    chrome = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
    if slides_pdf_dir.exists():
        for old in slides_pdf_dir.glob("slide_*.pdf"):
            old.unlink()
    slides_pdf_dir.mkdir(parents=True, exist_ok=True)

    width_px, height_px = 1920, 1080
    slide_pdfs: list[Path] = []

    with sync_playwright() as p:
        browser = p.chromium.launch(
            executable_path=chrome,
            headless=True,
            args=["--disable-gpu"],
        )
        page = browser.new_page(viewport={"width": width_px, "height": height_px})
        page.goto(html_path.as_uri(), wait_until="networkidle")
        page.add_style_tag(content="""
          .slide { overflow: hidden !important; height: 100vh !important; }
          .nav, .help { display: none !important; }
        """)

        for i in range(1, len(SLIDES) + 1):
            page.evaluate(
                """(n) => {
                  document.querySelectorAll('.slide').forEach((s, idx) => {
                    s.classList.toggle('active', idx === n - 1);
                    s.style.display = idx === n - 1 ? 'block' : 'none';
                  });
                }""",
                i,
            )
            page.wait_for_timeout(400)
            png_bytes = page.screenshot(type="png", full_page=False)
            img = Image.open(BytesIO(png_bytes)).convert("RGB")

            out = slides_pdf_dir / f"slide_{i:02d}.pdf"
            c = canvas.Canvas(str(out), pagesize=(13.333 * 72, 7.5 * 72))
            c.drawImage(ImageReader(img), 0, 0, width=13.333 * 72, height=7.5 * 72)
            c.showPage()
            c.save()
            slide_pdfs.append(out)
        browser.close()

    writer = PdfWriter()
    for sp in slide_pdfs:
        reader = PdfReader(str(sp))
        writer.add_page(reader.pages[0])
    with open(pdf_path, "wb") as f:
        writer.write(f)
    return len(SLIDES)


def split_pdf(pdf_path: Path, out_dir: Path) -> int:
    """Legacy helper — per-slide PDFs are written directly during PDF build."""
    reader = PdfReader(str(pdf_path))
    return len(reader.pages)


def copy_deliverables(files: list[Path]) -> None:
    if not DELIVERABLES.parent.exists():
        return
    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    slides_pdf_src = OUT / "slides_pdf"
    slides_pdf_dst = DELIVERABLES / "slides_pdf"
    for f in files:
        shutil.copy2(f, DELIVERABLES / f.name)
    if slides_pdf_src.exists():
        if slides_pdf_dst.exists():
            shutil.rmtree(slides_pdf_dst)
        shutil.copytree(slides_pdf_src, slides_pdf_dst)


def fmt_size(path: Path) -> str:
    n = path.stat().st_size
    if n < 1024:
        return f"{n} B"
    if n < 1024 * 1024:
        return f"{n / 1024:.1f} KB"
    return f"{n / (1024 * 1024):.1f} MB"


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)

    md_path = OUT / "Preflight_Sales_Deck_Slides.md"
    html_path = OUT / "Preflight_Sales_Deck.html"
    readme_path = OUT / "README.txt"
    pptx_path = OUT / "Preflight_Sales_Deck.pptx"
    pdf_path = OUT / "Preflight_Sales_Deck.pdf"
    slides_pdf_dir = OUT / "slides_pdf"

    md_path.write_text(build_markdown(), encoding="utf-8")
    html_path.write_text(build_html(), encoding="utf-8")
    readme_path.write_text(build_readme(), encoding="utf-8")
    build_pptx(pptx_path)
    page_count = build_pdf_from_html(html_path, pdf_path, slides_pdf_dir)

    deliverable_files = [md_path, html_path, readme_path, pptx_path, pdf_path]
    copy_deliverables(deliverable_files)

    print(f"Slide count: {len(SLIDES)} (PDF pages: {page_count})")
    print(f"Target: {OUT}")
    for p in deliverable_files:
        print(f"  {p.name}: {fmt_size(p)}")
    print(f"  slides_pdf/: {page_count} files")
    for p in sorted(slides_pdf_dir.glob("slide_*.pdf")):
        print(f"    {p.name}: {fmt_size(p)}")
    if DELIVERABLES.exists():
        print(f"Copied to: {DELIVERABLES}")


if __name__ == "__main__":
    main()
